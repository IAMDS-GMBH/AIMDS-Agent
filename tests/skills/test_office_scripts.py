from __future__ import annotations

import importlib.util
import json
import sys
from pathlib import Path

import pytest


REPO_ROOT = Path(__file__).resolve().parents[2]
PPT_SCRIPTS = REPO_ROOT / "skills/productivity/powerpoint/scripts"


def _load_module(path: Path, name: str):
    spec = importlib.util.spec_from_file_location(name, path)
    assert spec and spec.loader
    module = importlib.util.module_from_spec(spec)
    spec.loader.exec_module(module)
    return module


def test_word_convert_run_pandoc_handles_missing_binary(monkeypatch: pytest.MonkeyPatch) -> None:
    mod = _load_module(
        REPO_ROOT / "skills/productivity/word/scripts/convert.py",
        "word_convert_script",
    )

    def _raise_file_not_found(*args, **kwargs):  # type: ignore[no-untyped-def]
        raise FileNotFoundError("pandoc not found")

    monkeypatch.setattr(mod.subprocess, "run", _raise_file_not_found)
    ok, err = mod._run_pandoc(Path("in.docx"), Path("out.md"))
    assert ok is False
    assert err == "pandoc not found"


def test_excel_write_resolves_soffice_binary(monkeypatch: pytest.MonkeyPatch) -> None:
    mod = _load_module(
        REPO_ROOT / "skills/productivity/excel/scripts/write.py",
        "excel_write_script",
    )

    def _which(name: str) -> str | None:
        return "/usr/bin/soffice" if name == "soffice" else None

    monkeypatch.setattr(mod.shutil, "which", _which)
    assert mod._libreoffice_cmd() == "/usr/bin/soffice"


def test_excel_write_coerce_keeps_formulas_and_parses_scalars() -> None:
    mod = _load_module(REPO_ROOT / "skills/productivity/excel/scripts/write.py", "excel_write_coerce")
    assert mod._coerce("=SUM(A1:A2)") == "=SUM(A1:A2)"
    assert mod._coerce("42") == 42
    assert mod._coerce("4.5") == 4.5
    assert mod._coerce("true") is True
    assert mod._coerce("1,300") == "1,300"
    assert mod._coerce(7) == 7


def test_excel_write_parse_cells_accepts_dict_and_pairs() -> None:
    mod = _load_module(REPO_ROOT / "skills/productivity/excel/scripts/write.py", "excel_write_cells")
    assert mod._parse_cells('{"A1": 1, "B2": "x"}') == [("A1", 1), ("B2", "x")]
    assert mod._parse_cells('[["A1", 1], {"cell": "B2", "value": "x"}]') == [("A1", 1), ("B2", "x")]
    with pytest.raises(ValueError):
        mod._parse_cells('"A1"')


def test_excel_read_csv_single_cell_exits() -> None:
    mod = _load_module(
        REPO_ROOT / "skills/productivity/excel/scripts/read.py",
        "excel_read_script",
    )

    assert mod._is_csv("report.csv") is True
    assert mod._is_csv("report.xlsx") is False
    with pytest.raises(SystemExit):
        mod.read_cell("report.csv", "A1")


def test_word_replace_in_paragraph_spans_runs_and_keeps_other_formatting() -> None:
    pytest.importorskip("docx")
    from docx import Document

    mod = _load_module(REPO_ROOT / "skills/productivity/word/scripts/write.py", "word_write_script")
    doc = Document()
    para = doc.add_paragraph()
    para.add_run("Keep ")
    para.add_run("Ol")
    mid = para.add_run("d Te")
    mid.bold = True
    para.add_run("xt and Old Text again ")
    tail = para.add_run("tail")
    tail.italic = True

    assert mod._replace_in_paragraph(para, "Old Text", "New") == 2
    assert para.text == "Keep New and New again tail"
    assert para.runs[0].text == "Keep "
    assert para.runs[-1].text == "tail" and para.runs[-1].italic is True
    assert mod._replace_in_paragraph(para, "missing", "x") == 0


def test_deck_normalize_bullets_supports_levels() -> None:
    pytest.importorskip("pptx")
    mod = _load_module(PPT_SCRIPTS / "deck.py", "deck_script")
    items = mod._normalize_bullets(["top", {"text": "sub", "level": 1}, "  - indented", "", {"text": "  "}])
    assert items == [("top", 0), ("sub", 1), ("indented", 1)]


def test_deck_cli_roundtrip(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    mod = _load_module(PPT_SCRIPTS / "deck.py", "deck_cli")
    out = tmp_path / "deck.pptx"
    mod.main(["create", str(out), "--title", "T", "--slides-json", json.dumps([{"title": "A", "bullets": ["1"]}])])
    mod.main(["add-slide", str(out), "--title", "B", "--bullets-json", '["x"]', "--index", "2"])
    mod.main(["replace", str(out), "--find", "x", "--replace", "y"])
    mod.main(["delete-slide", str(out), "--index", "3"])

    from pptx import Presentation

    prs = Presentation(str(out))
    titles = [s.shapes.title.text for s in prs.slides]
    assert titles == ["T", "B"]
    with pytest.raises(SystemExit):
        mod.main(["delete-slide", str(out), "--index", "9"])


def test_unpack_pack_roundtrip_without_validators(tmp_path: Path) -> None:
    pytest.importorskip("pptx")
    pytest.importorskip("defusedxml")
    from pptx import Presentation

    src = tmp_path / "src.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Hello"
    prs.save(str(src))

    unpack = _load_module(PPT_SCRIPTS / "office/unpack.py", "office_unpack")
    pack = _load_module(PPT_SCRIPTS / "office/pack.py", "office_pack")
    assert pack.DOCXSchemaValidator is None  # validators module is not shipped

    unpacked = tmp_path / "unpacked"
    assert unpack.unpack(str(src), str(unpacked)).startswith("Unpacked")
    assert (unpacked / "ppt" / "slides" / "slide1.xml").exists()

    out = tmp_path / "out.pptx"
    _, message = pack.pack(str(unpacked), str(out), original_file=str(src), validate=True)
    assert message.startswith("Successfully packed"), message
    assert Presentation(str(out)).slides[0].shapes.title.text == "Hello"
