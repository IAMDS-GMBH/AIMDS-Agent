#!/usr/bin/env python3
"""Office productivity tool wrappers (Word/Excel/PowerPoint skill scripts)."""

from __future__ import annotations

import re
import os
import subprocess
import sys
import tempfile
from urllib.parse import unquote
from datetime import date
from pathlib import Path
from typing import Any

from tools.registry import registry, tool_error, tool_result


_REPO_ROOT = Path(__file__).resolve().parent.parent

_WORD_SCRIPT_ROOT = _REPO_ROOT / "skills" / "productivity" / "word" / "scripts"
_EXCEL_SCRIPT_ROOT = _REPO_ROOT / "skills" / "productivity" / "excel" / "scripts"
_PPT_SCRIPT_ROOT = _REPO_ROOT / "skills" / "productivity" / "powerpoint" / "scripts"

_DEFAULT_KPI_SHEETS = ["raw", "KPIs", "forecasts", "assumptions"]


def _check_office_tools() -> bool:
    return (
        (_WORD_SCRIPT_ROOT / "read.py").exists()
        and (_EXCEL_SCRIPT_ROOT / "read.py").exists()
        and (_PPT_SCRIPT_ROOT / "add_slide.py").exists()
    )


def _safe_str(value: Any) -> str | None:
    if value is None:
        return None
    text = str(value).strip()
    return text or None


def _safe_int(value: Any, default: int) -> int:
    try:
        return int(value)
    except (TypeError, ValueError):
        return default


def _workspace_root() -> Path:
    raw = (os.environ.get("TERMINAL_CWD") or "").strip()
    if raw:
        candidate = Path(raw).expanduser()
        if candidate.is_absolute() and candidate.is_dir():
            return candidate
    return Path.cwd()


def _resolve_workspace_path(path_value: str) -> str:
    p = Path(path_value).expanduser()
    if p.is_absolute():
        return str(p)
    return str((_workspace_root() / p).resolve())


def _materialize_markdown_source(source_path: str) -> str:
    prefix = "data:text/markdown,"
    if source_path.startswith(prefix):
        markdown_text = unquote(source_path[len(prefix):])
        tmp = tempfile.NamedTemporaryFile(
            mode="w",
            suffix=".md",
            prefix="office_word_",
            dir=str(_workspace_root()),
            encoding="utf-8",
            delete=False,
        )
        try:
            tmp.write(markdown_text)
            return tmp.name
        finally:
            tmp.close()
    return _resolve_workspace_path(source_path)


def _normalize_action(action: str | None) -> str:
    if not action:
        return ""
    return action.strip().lower().replace("-", "_").replace(" ", "_")


def _run_script(script_path: Path, args: list[str]) -> str:
    if not script_path.exists():
        return tool_error(f"Missing script: {script_path}")

    cmd = [sys.executable, str(script_path), *args]
    try:
        completed = subprocess.run(
            cmd,
            capture_output=True,
            text=True,
            cwd=str(_REPO_ROOT),
        )
    except Exception as exc:
        return tool_error(f"Failed to execute {script_path.name}: {exc}")

    ok = completed.returncode == 0
    return tool_result(
        success=ok,
        command=cmd,
        returncode=completed.returncode,
        stdout=(completed.stdout or "").strip(),
        stderr=(completed.stderr or "").strip(),
    )


def _month_labels(months: int) -> list[str]:
    months = max(1, min(24, months))
    today = date.today()
    labels: list[str] = []
    year = today.year
    month = today.month
    for idx in range(months):
        m = month - (months - 1 - idx)
        y = year
        while m <= 0:
            m += 12
            y -= 1
        labels.append(f"{y}-{m:02d}")
    return labels


def _markdown_to_slide_sections(markdown: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_bullets: list[str] = []

    for raw in markdown.splitlines():
        line = raw.strip()
        if not line:
            continue
        if line.startswith("#"):
            if current_title and current_bullets:
                sections.append((current_title, current_bullets[:8]))
            current_title = line.lstrip("#").strip() or "Section"
            current_bullets = []
            continue
        if line.startswith(("- ", "* ")):
            current_bullets.append(line[2:].strip())
        elif re.match(r"^\d+\.\s+", line):
            current_bullets.append(re.sub(r"^\d+\.\s+", "", line))
        elif current_title:
            current_bullets.append(line)

    if current_title and current_bullets:
        sections.append((current_title, current_bullets[:8]))

    if sections:
        return sections[:12]

    plain_lines = [line.strip() for line in markdown.splitlines() if line.strip()]
    if not plain_lines:
        return [("Summary", ["No content provided"])]
    return [("Summary", plain_lines[:10])]


def _validate_workbook(path: Path, required_sheets: list[str]) -> dict[str, Any]:
    try:
        import openpyxl
    except ImportError:
        return {
            "success": False,
            "error": "openpyxl is required for workbook validation",
        }

    if not path.exists():
        return {"success": False, "error": f"Workbook does not exist: {path}"}

    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    try:
        missing = [s for s in required_sheets if s not in wb.sheetnames]
        if missing:
            return {
                "success": False,
                "error": "Workbook is missing required sheets",
                "missing_sheets": missing,
                "sheets": wb.sheetnames,
            }
        non_empty = {
            s: bool((wb[s].max_row or 0) > 1 or (wb[s].max_column or 0) > 1)
            for s in required_sheets
        }
        return {
            "success": True,
            "sheets": wb.sheetnames,
            "required_sheets": required_sheets,
            "non_empty_checks": non_empty,
        }
    finally:
        wb.close()


def _generate_kpi_workbook(output_path: Path, months: int) -> dict[str, Any]:
    try:
        import openpyxl
    except ImportError:
        return {
            "success": False,
            "error": "openpyxl is required for generate_kpi_workbook",
        }

    labels = _month_labels(months)
    wb = openpyxl.Workbook()
    wb.remove(wb.active)

    ws_raw = wb.create_sheet("raw")
    ws_raw.append(
        [
            "month",
            "revenue",
            "churn_pct",
            "cac",
            "nps",
            "sla_pct",
            "outages",
            "cloud_cost",
        ]
    )
    for idx, label in enumerate(labels):
        revenue = 900_000 + (idx * 18_000) + (idx % 3) * 7_500
        churn = round(4.1 - (idx * 0.08), 2)
        cac = 820 + (idx % 4) * 25
        nps = 38 + min(12, idx)
        sla = round(99.2 + (idx % 4) * 0.08, 2)
        outages = max(0, 4 - (idx // 4))
        cloud = 210_000 + idx * 5_500
        ws_raw.append([label, revenue, churn, cac, nps, sla, outages, cloud])

    ws_kpi = wb.create_sheet("KPIs")
    ws_kpi.append(["metric", "value"])
    ws_kpi.append(["avg_revenue", f"=AVERAGE(raw!B2:B{months+1})"])
    ws_kpi.append(["latest_revenue", f"=raw!B{months+1}"])
    ws_kpi.append(["avg_churn_pct", f"=AVERAGE(raw!C2:C{months+1})"])
    ws_kpi.append(["avg_nps", f"=AVERAGE(raw!E2:E{months+1})"])
    ws_kpi.append(["avg_sla_pct", f"=AVERAGE(raw!F2:F{months+1})"])
    ws_kpi.append(["total_outages", f"=SUM(raw!G2:G{months+1})"])

    ws_fc = wb.create_sheet("forecasts")
    ws_fc.append(["scenario", "month", "revenue", "cloud_cost"])
    base_revenue = 900_000 + (months * 18_000)
    base_cloud = 210_000 + months * 5_500
    for scenario, rev_factor, cloud_factor in [
        ("base", 1.0, 1.0),
        ("upside", 1.08, 0.96),
        ("downside", 0.9, 1.12),
    ]:
        for step in range(1, 4):
            ws_fc.append(
                [
                    scenario,
                    f"M+{step}",
                    int(base_revenue * rev_factor * (1 + (step * 0.01))),
                    int(base_cloud * cloud_factor * (1 + (step * 0.015))),
                ]
            )

    ws_ass = wb.create_sheet("assumptions")
    ws_ass.append(["key", "value"])
    ws_ass.append(["report_months", months])
    ws_ass.append(["revenue_growth_rate", "1.8% / month"])
    ws_ass.append(["cost_growth_rate", "1.5% / month"])
    ws_ass.append(["scenario_upside_revenue_factor", 1.08])
    ws_ass.append(["scenario_downside_revenue_factor", 0.9])

    output_path.parent.mkdir(parents=True, exist_ok=True)
    wb.save(output_path)

    checks = _validate_workbook(output_path, _DEFAULT_KPI_SHEETS)
    checks["file_size_bytes"] = output_path.stat().st_size if output_path.exists() else 0
    return checks


def _validate_word_report(path: Path) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError:
        return {"success": False, "error": "python-docx is required for report validation"}

    if not path.exists():
        return {"success": False, "error": f"Report does not exist: {path}"}

    doc = Document(path)
    headings = [p.text.strip() for p in doc.paragraphs if p.style.name.startswith("Heading")]
    required = [
        "Executive Summary",
        "KPI Highlights",
        "Scenario Analysis",
        "Risk Register",
        "Action Plan",
        "Unknowns & Assumptions",
    ]
    missing = [h for h in required if h not in headings]
    non_empty_paragraphs = sum(1 for p in doc.paragraphs if p.text.strip())
    return {
        "success": len(missing) == 0 and non_empty_paragraphs >= 10,
        "headings": headings,
        "required_headings": required,
        "missing_headings": missing,
        "non_empty_paragraphs": non_empty_paragraphs,
        "file_size_bytes": path.stat().st_size,
    }


def _generate_exec_report(output_path: Path, title: str) -> dict[str, Any]:
    try:
        from docx import Document
    except ImportError:
        return {"success": False, "error": "python-docx is required for generate_exec_report"}

    doc = Document()
    doc.add_heading(title or "Q3 Operational Review", level=1)

    sections = {
        "Executive Summary": [
            "This report summarizes operational performance over the last 12 months.",
            "Overall trend shows improving service quality with controlled cost growth.",
        ],
        "KPI Highlights": [
            "Revenue trend remains positive with moderate month-over-month growth.",
            "SLA performance stayed above target while outages declined quarter-over-quarter.",
        ],
        "Scenario Analysis": [
            "Base: maintain current growth and cost control trajectory.",
            "Upside: faster topline acceleration with improved acquisition efficiency.",
            "Downside: demand contraction with elevated infrastructure pressure.",
        ],
        "Risk Register": [
            "Capacity saturation risk in peak workload windows.",
            "Vendor concentration risk in core cloud services.",
            "Execution risk on cross-team delivery commitments.",
        ],
        "Action Plan": [
            "Finalize cost anomaly alerts and ownership SLAs.",
            "Harden incident runbooks for top recurring failure modes.",
            "Prioritize CAC optimization experiments by segment.",
        ],
        "Unknowns & Assumptions": [
            "Assumes no disruptive pricing changes from major vendors.",
            "Assumes current support staffing and on-call cadence remain stable.",
        ],
    }
    for heading, bullets in sections.items():
        doc.add_heading(heading, level=2)
        for line in bullets:
            doc.add_paragraph(line, style="List Bullet")

    output_path.parent.mkdir(parents=True, exist_ok=True)
    doc.save(output_path)
    return _validate_word_report(output_path)


def _validate_pptx(path: Path, min_slides: int = 10, min_charts: int = 3) -> dict[str, Any]:
    try:
        from pptx import Presentation
    except ImportError:
        return {"success": False, "error": "python-pptx is required for deck validation"}

    if not path.exists():
        return {"success": False, "error": f"Deck does not exist: {path}"}

    prs = Presentation(path)
    slide_count = len(prs.slides)
    chart_count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if getattr(shape, "has_chart", False):
                chart_count += 1
    return {
        "success": slide_count >= min_slides and chart_count >= min_charts,
        "slides": slide_count,
        "charts": chart_count,
        "min_slides": min_slides,
        "min_charts": min_charts,
        "file_size_bytes": path.stat().st_size,
    }


def _generate_review_deck(output_path: Path, title: str, content_markdown: str | None = None) -> dict[str, Any]:
    try:
        from pptx import Presentation
        from pptx.chart.data import CategoryChartData
        from pptx.enum.chart import XL_CHART_TYPE
        from pptx.util import Inches
    except ImportError:
        return {"success": False, "error": "python-pptx is required for generate_review_deck"}

    prs = Presentation()
    title_text = title or "Q3 Operational Review Pack"

    # Slide 1
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title_text
    slide.placeholders[1].text = "Generated from provided input content"

    if content_markdown and content_markdown.strip():
        for heading, bullets in _markdown_to_slide_sections(content_markdown):
            s = prs.slides.add_slide(prs.slide_layouts[1])
            s.shapes.title.text = heading
            body = s.shapes.placeholders[1].text_frame
            body.clear()
            for idx, bullet in enumerate(bullets):
                p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
                p.text = bullet
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(output_path)
        return _validate_pptx(output_path, min_slides=2, min_charts=0)

    # Slides 2-4 with charts
    chart_specs = [
        ("Revenue trend", [980, 1020, 1055, 1110, 1160, 1215]),
        ("Cloud cost trend", [220, 229, 235, 244, 252, 261]),
        ("NPS trend", [39, 42, 44, 46, 47, 49]),
    ]
    months = ["Jan", "Feb", "Mar", "Apr", "May", "Jun"]
    for title_line, values in chart_specs:
        s = prs.slides.add_slide(prs.slide_layouts[5])
        s.shapes.title.text = title_line
        data = CategoryChartData()
        data.categories = months
        data.add_series(title_line, values)
        s.shapes.add_chart(
            XL_CHART_TYPE.LINE_MARKERS,
            Inches(0.8),
            Inches(1.6),
            Inches(8.5),
            Inches(4.2),
            data,
        )

    # Slides 5-10 content
    body_slides = [
        ("KPI Highlights", ["Revenue up 18%", "SLA 99.4%", "Outages down 35%"]),
        ("Scenario Analysis", ["Base: stable growth", "Upside: +8% acceleration", "Downside: -10% demand"]),
        ("Risk Register", ["Cloud concentration", "On-call fatigue", "Release coupling"]),
        ("Action Plan", ["Cost guardrails", "Runbook hardening", "Quarterly experiment cadence"]),
        ("Unknowns & Assumptions", ["No major price shocks", "Staffing remains constant", "Demand seasonality holds"]),
        ("Decision Summary", ["Proceed with base plan", "Prepare downside response", "Track leading indicators weekly"]),
    ]
    for heading, bullets in body_slides:
        s = prs.slides.add_slide(prs.slide_layouts[1])
        s.shapes.title.text = heading
        body = s.shapes.placeholders[1].text_frame
        body.clear()
        for idx, bullet in enumerate(bullets):
            p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
            p.text = bullet

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return _validate_pptx(output_path)


def office_word_tool(args: dict, **kwargs) -> str:
    action = _normalize_action(_safe_str(args.get("action")))
    if not action:
        return tool_error("action is required")

    path = _safe_str(args.get("path"))
    output_path = _safe_str(args.get("output_path"))
    title = _safe_str(args.get("title"))
    text = _safe_str(args.get("text"))
    find_text = _safe_str(args.get("find_text"))
    replace_text = _safe_str(args.get("replace_text"))
    source_path = _safe_str(args.get("source_path"))
    source_paths = args.get("source_paths") or []
    fmt = (_safe_str(args.get("format")) or "").lower()

    action_aliases = {
        "read": "read_text",
        "read_doc": "read_text",
        "read_docx": "read_text",
        "markdown": "read_markdown",
        "read_md": "read_markdown",
        "tables": "read_tables",
        "metadata": "read_metadata",
        "styles": "read_styles",
        "from_md": "from_markdown",
        "replace": "find_replace",
        "append": "append_paragraph",
        "merge_docs": "merge",
        "convert_pdf": "convert",
        "convert_md": "convert",
        "convert_html": "convert",
        "generate_report": "generate_exec_report",
        "validate": "validate_report",
    }
    action = action_aliases.get(action, action)

    if action in {"read_text", "read_markdown", "read_tables", "read_metadata", "read_styles"}:
        if not path:
            return tool_error("path is required for read actions")
        action_args = [path]
        flag_map = {
            "read_markdown": "--markdown",
            "read_tables": "--tables",
            "read_metadata": "--metadata",
            "read_styles": "--styles",
        }
        flag = flag_map.get(action)
        if flag:
            action_args.append(flag)
        return _run_script(_WORD_SCRIPT_ROOT / "read.py", action_args)

    if action == "from_markdown":
        if not source_path or not output_path:
            return tool_error("source_path and output_path are required for from_markdown")
        resolved_source = _materialize_markdown_source(source_path)
        resolved_output = _resolve_workspace_path(output_path)
        return _run_script(_WORD_SCRIPT_ROOT / "write.py", ["--from-md", resolved_source, resolved_output])

    if action == "find_replace":
        if not path or find_text is None or replace_text is None:
            return tool_error("path, find_text, and replace_text are required for find_replace")
        return _run_script(_WORD_SCRIPT_ROOT / "write.py", ["--replace", find_text, replace_text, path])

    if action == "append_paragraph":
        if not path or text is None:
            return tool_error("path and text are required for append_paragraph")
        return _run_script(_WORD_SCRIPT_ROOT / "write.py", ["--append", text, path])

    if action == "merge":
        if not source_paths or not isinstance(source_paths, list) or not output_path:
            return tool_error("source_paths (array) and output_path are required for merge")
        paths = [_resolve_workspace_path(str(p)) for p in source_paths if _safe_str(p)]
        if len(paths) < 2:
            return tool_error("merge requires at least two source_paths")
        return _run_script(
            _WORD_SCRIPT_ROOT / "write.py",
            ["--merge", *paths, "--out", _resolve_workspace_path(output_path)],
        )

    if action == "convert":
        if not path or fmt not in {"pdf", "md", "markdown", "html", "htm"}:
            return tool_error("convert requires path and format in [pdf, md, markdown, html, htm]")
        return _run_script(_WORD_SCRIPT_ROOT / "convert.py", [path, "--to", fmt])

    if action == "generate_exec_report":
        if not output_path:
            return tool_error("output_path is required for generate_exec_report")
        checks = _generate_exec_report(Path(output_path), title or "Q3 Operational Review")
        return tool_result(action=action, output_path=output_path, checks=checks, success=bool(checks.get("success")))

    if action == "validate_report":
        if not path:
            return tool_error("path is required for validate_report")
        checks = _validate_word_report(Path(path))
        return tool_result(action=action, path=path, checks=checks, success=bool(checks.get("success")))

    return tool_error(
        "Unsupported word action",
        supported_actions=[
            "read_text",
            "read_markdown",
            "read_tables",
            "read_metadata",
            "read_styles",
            "from_markdown",
            "find_replace",
            "append_paragraph",
            "merge",
            "convert",
            "generate_exec_report",
            "validate_report",
        ],
    )


def office_excel_tool(args: dict, **kwargs) -> str:
    action = _normalize_action(_safe_str(args.get("action")))
    if not action:
        return tool_error("action is required")

    path = _safe_str(args.get("path"))
    output_path = _safe_str(args.get("output_path"))
    sheet = _safe_str(args.get("sheet"))
    cell = _safe_str(args.get("cell"))
    value = args.get("value")
    row_csv = _safe_str(args.get("row_csv"))
    source_path = _safe_str(args.get("source_path"))
    months = _safe_int(args.get("months"), 12)
    required_sheets = args.get("required_sheets")

    action_aliases = {
        "read": "read_sheet",
        "list": "list_sheets",
        "sheets": "list_sheets",
        "statistics": "stats",
        "describe": "stats",
        "cell": "read_cell",
        "meta": "metadata",
        "csv_to_xlsx": "from_csv",
        "fromcsv": "from_csv",
        "set": "set_cell",
        "append": "append_row",
        "export_csv": "to_csv",
        "export_pdf": "to_pdf",
        "pdf": "to_pdf",
        "generate_workbook": "generate_kpi_workbook",
        "validate": "validate_workbook",
    }
    action = action_aliases.get(action, action)

    if action in {"read_sheet", "list_sheets", "stats", "read_cell", "metadata"}:
        if not path:
            return tool_error("path is required for read actions")
        cmd = [path]
        if action == "list_sheets":
            cmd.append("--sheets")
        elif action == "stats":
            cmd.append("--stats")
        elif action == "read_cell":
            if not cell:
                return tool_error("cell is required for read_cell")
            cmd.extend(["--cell", cell])
        elif action == "metadata":
            cmd.append("--metadata")
        if sheet:
            cmd.extend(["--sheet", sheet])
        return _run_script(_EXCEL_SCRIPT_ROOT / "read.py", cmd)

    if action == "from_csv":
        if not source_path or not output_path:
            return tool_error("source_path and output_path are required for from_csv")
        return _run_script(_EXCEL_SCRIPT_ROOT / "write.py", ["--from-csv", source_path, output_path])

    if action == "set_cell":
        if not path or not cell or value is None:
            return tool_error("path, cell, and value are required for set_cell")
        cmd = ["--set-cell", cell, str(value), path]
        if sheet:
            cmd.extend(["--sheet", sheet])
        return _run_script(_EXCEL_SCRIPT_ROOT / "write.py", cmd)

    if action == "append_row":
        if not path or not row_csv:
            return tool_error("path and row_csv are required for append_row")
        cmd = ["--append-row", row_csv, path]
        if sheet:
            cmd.extend(["--sheet", sheet])
        return _run_script(_EXCEL_SCRIPT_ROOT / "write.py", cmd)

    if action == "to_csv":
        if not path:
            return tool_error("path is required for to_csv")
        cmd = ["--to-csv", path]
        if output_path:
            cmd.append(output_path)
        if sheet:
            cmd.extend(["--sheet", sheet])
        return _run_script(_EXCEL_SCRIPT_ROOT / "write.py", cmd)

    if action == "to_pdf":
        if not path:
            return tool_error("path is required for to_pdf")
        return _run_script(_EXCEL_SCRIPT_ROOT / "write.py", ["--to-pdf", path])

    if action == "generate_kpi_workbook":
        if not output_path:
            return tool_error("output_path is required for generate_kpi_workbook")
        checks = _generate_kpi_workbook(Path(output_path), months)
        return tool_result(
            action=action,
            output_path=output_path,
            months=months,
            checks=checks,
            success=bool(checks.get("success")),
        )

    if action == "validate_workbook":
        if not path:
            return tool_error("path is required for validate_workbook")
        sheet_list = required_sheets if isinstance(required_sheets, list) and required_sheets else _DEFAULT_KPI_SHEETS
        checks = _validate_workbook(Path(path), [str(s) for s in sheet_list])
        return tool_result(action=action, path=path, checks=checks, success=bool(checks.get("success")))

    return tool_error(
        "Unsupported excel action",
        supported_actions=[
            "read_sheet",
            "list_sheets",
            "stats",
            "read_cell",
            "metadata",
            "from_csv",
            "set_cell",
            "append_row",
            "to_csv",
            "to_pdf",
            "generate_kpi_workbook",
            "validate_workbook",
        ],
    )


def office_powerpoint_tool(args: dict, **kwargs) -> str:
    action = _normalize_action(_safe_str(args.get("action")))
    if not action:
        return tool_error("action is required")

    unpacked_dir = _safe_str(args.get("unpacked_dir"))
    source = _safe_str(args.get("source"))
    input_directory = _safe_str(args.get("input_directory"))
    output_file = _safe_str(args.get("output_file"))
    source_path = _safe_str(args.get("source_path"))
    text = _safe_str(args.get("text"))
    original_file = _safe_str(args.get("original_file"))
    title = _safe_str(args.get("title"))
    validate = args.get("validate")
    min_slides = _safe_int(args.get("min_slides"), 10)
    min_charts = _safe_int(args.get("min_charts"), 3)

    action_aliases = {
        "add": "add_slide",
        "duplicate_slide": "add_slide",
        "cleanup": "clean",
        "repack": "pack",
        "generate_deck": "generate_review_deck",
        "validate": "validate_deck",
    }
    action = action_aliases.get(action, action)

    if action == "add_slide":
        if not unpacked_dir or not source:
            return tool_error("unpacked_dir and source are required for add_slide")
        return _run_script(_PPT_SCRIPT_ROOT / "add_slide.py", [unpacked_dir, source])

    if action == "clean":
        if not unpacked_dir:
            return tool_error("unpacked_dir is required for clean")
        return _run_script(_PPT_SCRIPT_ROOT / "clean.py", [unpacked_dir])

    if action == "pack":
        if not input_directory or not output_file:
            return tool_error("input_directory and output_file are required for pack")
        cmd = [input_directory, output_file]
        if original_file:
            cmd.extend(["--original", original_file])
        if isinstance(validate, bool):
            cmd.extend(["--validate", "true" if validate else "false"])
        return _run_script(_PPT_SCRIPT_ROOT / "office" / "pack.py", cmd)

    if action == "generate_review_deck":
        if not output_file:
            return tool_error("output_file is required for generate_review_deck")
        content_markdown: str | None = text
        if source_path:
            try:
                content_markdown = Path(_resolve_workspace_path(source_path)).read_text(encoding="utf-8")
            except OSError as exc:
                return tool_error(f"Failed to read source_path for generate_review_deck: {exc}")
        checks = _generate_review_deck(Path(_resolve_workspace_path(output_file)), title or "Q3 Operational Review Pack", content_markdown=content_markdown)
        return tool_result(action=action, output_file=output_file, checks=checks, success=bool(checks.get("success")))

    if action == "validate_deck":
        deck_path = output_file or _safe_str(args.get("path"))
        if not deck_path:
            return tool_error("output_file or path is required for validate_deck")
        checks = _validate_pptx(Path(deck_path), min_slides=min_slides, min_charts=min_charts)
        return tool_result(action=action, path=deck_path, checks=checks, success=bool(checks.get("success")))

    return tool_error(
        "Unsupported powerpoint action",
        supported_actions=["add_slide", "clean", "pack", "generate_review_deck", "validate_deck"],
    )


OFFICE_WORD_SCHEMA = {
    "name": "office_word",
    "description": (
        "Operate on Word files via deterministic skill wrappers. "
        "Action-specific required args: "
        "read_* -> path; from_markdown -> source_path+output_path; "
        "find_replace -> path+find_text+replace_text; append_paragraph -> path+text; "
        "merge -> source_paths+output_path; convert -> path+format; "
        "generate_exec_report -> output_path; validate_report -> path."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": [
                    "read_text",
                    "read_markdown",
                    "read_tables",
                    "read_metadata",
                    "read_styles",
                    "from_markdown",
                    "find_replace",
                    "append_paragraph",
                    "merge",
                    "convert",
                    "generate_exec_report",
                    "validate_report",
                ],
            },
            "path": {"type": "string"},
            "source_path": {"type": "string"},
            "source_paths": {"type": "array", "items": {"type": "string"}},
            "output_path": {"type": "string"},
            "title": {"type": "string"},
            "format": {"type": "string", "enum": ["pdf", "md", "markdown", "html", "htm"]},
            "text": {"type": "string"},
            "find_text": {"type": "string"},
            "replace_text": {"type": "string"},
        },
        "required": ["action"],
        "oneOf": [
            {
                "properties": {
                    "action": {
                        "enum": [
                            "read_text",
                            "read_markdown",
                            "read_tables",
                            "read_metadata",
                            "read_styles",
                        ]
                    }
                },
                "required": ["action", "path"],
            },
            {
                "properties": {"action": {"const": "from_markdown"}},
                "required": ["action", "source_path", "output_path"],
            },
            {
                "properties": {"action": {"const": "find_replace"}},
                "required": ["action", "path", "find_text", "replace_text"],
            },
            {
                "properties": {"action": {"const": "append_paragraph"}},
                "required": ["action", "path", "text"],
            },
            {
                "properties": {"action": {"const": "merge"}},
                "required": ["action", "source_paths", "output_path"],
            },
            {
                "properties": {"action": {"const": "convert"}},
                "required": ["action", "path", "format"],
            },
            {
                "properties": {"action": {"const": "generate_exec_report"}},
                "required": ["action", "output_path"],
            },
            {
                "properties": {"action": {"const": "validate_report"}},
                "required": ["action", "path"],
            },
        ],
        "additionalProperties": False,
    },
}

OFFICE_EXCEL_SCHEMA = {
    "name": "office_excel",
    "description": "Operate on Excel/CSV files via deterministic skill wrappers.",
    "parameters": {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": [
                    "read_sheet",
                    "list_sheets",
                    "stats",
                    "read_cell",
                    "metadata",
                    "from_csv",
                    "set_cell",
                    "append_row",
                    "to_csv",
                    "to_pdf",
                    "generate_kpi_workbook",
                    "validate_workbook",
                ],
            },
            "path": {"type": "string"},
            "source_path": {"type": "string"},
            "output_path": {"type": "string"},
            "sheet": {"type": "string"},
            "cell": {"type": "string"},
            "value": {"type": ["string", "number", "boolean"]},
            "row_csv": {"type": "string"},
            "months": {"type": "integer", "minimum": 1, "maximum": 24},
            "required_sheets": {"type": "array", "items": {"type": "string"}},
        },
        "required": ["action"],
        "additionalProperties": False,
    },
}

OFFICE_POWERPOINT_SCHEMA = {
    "name": "office_powerpoint",
    "description": (
        "Operate on PowerPoint artifacts via deterministic skill wrappers. "
        "Action-specific required args: add_slide -> unpacked_dir+source; "
        "clean -> unpacked_dir; pack -> input_directory+output_file; "
        "generate_review_deck -> output_file (+ optional source_path or text for dynamic content); "
        "validate_deck -> output_file or path."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "action": {
                "type": "string",
                "enum": [
                    "add_slide",
                    "clean",
                    "pack",
                    "generate_review_deck",
                    "validate_deck",
                ],
            },
            "unpacked_dir": {"type": "string"},
            "source": {"type": "string"},
            "input_directory": {"type": "string"},
            "output_file": {"type": "string"},
            "source_path": {"type": "string"},
            "text": {"type": "string"},
            "path": {"type": "string"},
            "original_file": {"type": "string"},
            "title": {"type": "string"},
            "validate": {"type": "boolean"},
            "min_slides": {"type": "integer", "minimum": 1},
            "min_charts": {"type": "integer", "minimum": 0},
        },
        "required": ["action"],
        "oneOf": [
            {
                "properties": {"action": {"const": "add_slide"}},
                "required": ["action", "unpacked_dir", "source"],
            },
            {
                "properties": {"action": {"const": "clean"}},
                "required": ["action", "unpacked_dir"],
            },
            {
                "properties": {"action": {"const": "pack"}},
                "required": ["action", "input_directory", "output_file"],
            },
            {
                "properties": {"action": {"const": "generate_review_deck"}},
                "required": ["action", "output_file"],
            },
            {
                "properties": {"action": {"const": "validate_deck"}},
                "required": ["action"],
                "anyOf": [
                    {"required": ["output_file"]},
                    {"required": ["path"]},
                ],
            },
        ],
        "additionalProperties": False,
    },
}


registry.register(
    name="office_word",
    toolset="office",
    schema=OFFICE_WORD_SCHEMA,
    handler=office_word_tool,
    check_fn=_check_office_tools,
    emoji="📝",
)

registry.register(
    name="office_excel",
    toolset="office",
    schema=OFFICE_EXCEL_SCHEMA,
    handler=office_excel_tool,
    check_fn=_check_office_tools,
    emoji="📊",
)

registry.register(
    name="office_powerpoint",
    toolset="office",
    schema=OFFICE_POWERPOINT_SCHEMA,
    handler=office_powerpoint_tool,
    check_fn=_check_office_tools,
    emoji="📽️",
)
