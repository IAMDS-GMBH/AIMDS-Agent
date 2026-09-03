#!/usr/bin/env python3
"""Create, inspect and edit PowerPoint (.pptx) decks directly with python-pptx.

No unpack/pack round trip: every command opens the .pptx, applies the change
and saves it again. Use the XML workflow (office/unpack.py, add_slide.py,
clean.py, office/pack.py) only for edits python-pptx cannot express.

Usage:
    python deck.py read deck.pptx [--json] [--markdown]
    python deck.py layouts deck.pptx
    python deck.py create out.pptx --title "Title" [--subtitle "Sub"]
                   [--slides-json '[{"title": "...", "bullets": ["a", {"text": "b", "level": 1}], "notes": "..."}]']
                   [--template base.pptx] [--keep-template-slides]
    python deck.py add-slide deck.pptx --title "T" [--bullets-json '["a","b"]'] [--notes "..."]
                   [--layout 1 | --layout "Title and Content"] [--index 2]
    python deck.py delete-slide deck.pptx --index 3
    python deck.py replace deck.pptx --find "Old" --replace "New"
    python deck.py to-pdf deck.pptx [--outdir DIR]

Slide indices are 1-based. ``--slides-json`` may also be a plain list of
strings (each becomes a title-only slide).
"""

from __future__ import annotations

import argparse
import json
import shutil
import subprocess
import sys
from pathlib import Path


# --------------------------------------------------------------------------- helpers

def _presentation(path: str | Path | None = None):
    from pptx import Presentation

    return Presentation(str(path)) if path else Presentation()


def _slide_title(slide) -> str:
    try:
        title = slide.shapes.title
    except Exception:
        title = None
    return title.text if title is not None else ""


def _iter_text_frames(shapes):
    """Yield every text frame under ``shapes`` (groups and tables included)."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    for shape in shapes:
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from _iter_text_frames(shape.shapes)
            continue
        if getattr(shape, "has_table", False) and shape.has_table:
            for row in shape.table.rows:
                for cell in row.cells:
                    yield cell.text_frame
            continue
        if shape.has_text_frame:
            yield shape.text_frame


def _paragraph_text(paragraph) -> str:
    return "".join(run.text for run in paragraph.runs) or paragraph.text


def _replace_in_paragraph(paragraph, old: str, new: str) -> int:
    """Run-spanning replace; keeps the first touched run's formatting."""
    runs = list(paragraph.runs)
    if not runs or not old:
        return 0
    texts = [r.text for r in runs]
    full = "".join(texts)
    if old not in full:
        return 0

    positions = []
    i = full.find(old)
    while i != -1:
        positions.append(i)
        i = full.find(old, i + len(old))

    bounds = []
    offset = 0
    for t in texts:
        bounds.append((offset, offset + len(t)))
        offset += len(t)

    def _run_at(pos: int, *, end: bool) -> int:
        for idx, (start, stop) in enumerate(bounds):
            if start == stop:
                continue
            if end:
                if start < pos <= stop:
                    return idx
            elif start <= pos < stop:
                return idx
        return len(bounds) - 1

    for pos in reversed(positions):
        match_end = pos + len(old)
        first = _run_at(pos, end=False)
        last = _run_at(match_end, end=True)
        a = pos - bounds[first][0]
        b = match_end - bounds[last][0]
        if first == last:
            texts[first] = texts[first][:a] + new + texts[first][b:]
        else:
            texts[first] = texts[first][:a] + new
            for mid in range(first + 1, last):
                texts[mid] = ""
            texts[last] = texts[last][b:]

    for run, text in zip(runs, texts):
        if run.text != text:
            run.text = text
    return len(positions)


def _pick_layout(prs, spec, *, default_names=("Title and Content",), default_index=1):
    layouts = list(prs.slide_layouts)
    if not layouts:
        raise SystemExit("Error: presentation has no slide layouts")
    if spec is None or spec == "":
        for name in default_names:
            for layout in layouts:
                if layout.name.strip().lower() == name.lower():
                    return layout
        return layouts[default_index] if len(layouts) > default_index else layouts[0]
    if isinstance(spec, int) or (isinstance(spec, str) and spec.strip().isdigit()):
        idx = int(spec)
        if 0 <= idx < len(layouts):
            return layouts[idx]
        raise SystemExit(f"Error: layout index {idx} out of range (0..{len(layouts) - 1})")
    wanted = str(spec).strip().lower()
    for layout in layouts:
        if layout.name.strip().lower() == wanted:
            return layout
    for layout in layouts:
        if wanted in layout.name.strip().lower():
            return layout
    names = [layout.name for layout in layouts]
    raise SystemExit(f"Error: layout {spec!r} not found. Available: {names}")


def _body_placeholder(slide):
    from pptx.enum.shapes import PP_PLACEHOLDER

    title_types = {PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE, PP_PLACEHOLDER.SUBTITLE}
    for shape in slide.placeholders:
        if shape.placeholder_format.type in title_types:
            continue
        if shape.has_text_frame:
            return shape
    return None


def _normalize_bullets(bullets) -> list[tuple[str, int]]:
    out: list[tuple[str, int]] = []
    for item in bullets or []:
        if isinstance(item, dict):
            text = str(item.get("text", "")).strip()
            level = int(item.get("level", 0) or 0)
        else:
            raw = str(item)
            stripped = raw.lstrip(" ")
            level = (len(raw) - len(stripped)) // 2
            text = stripped[2:].strip() if stripped.startswith(("- ", "* ")) else stripped.strip()
        if text:
            out.append((text, max(0, min(level, 4))))
    return out


def _fill_slide(prs, slide, title: str | None, bullets, notes: str | None) -> None:
    from pptx.util import Inches

    if title is not None and slide.shapes.title is not None:
        slide.shapes.title.text = title

    items = _normalize_bullets(bullets)
    if items:
        body = _body_placeholder(slide)
        if body is None:
            left = Inches(0.7)
            top = Inches(1.5)
            width = prs.slide_width - Inches(1.4)
            height = prs.slide_height - Inches(2.2)
            body = slide.shapes.add_textbox(left, top, width, height)
            body.text_frame.word_wrap = True
        tf = body.text_frame
        tf.clear()
        for idx, (text, level) in enumerate(items):
            paragraph = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            paragraph.text = text
            paragraph.level = level

    if notes:
        slide.notes_slide.notes_text_frame.text = notes


def _drop_slide(prs, sld_id) -> None:
    prs.part.drop_rel(sld_id.rId)
    prs.slides._sldIdLst.remove(sld_id)


def _summary(path, prs, **extra) -> None:
    payload = {"file": str(path), "slides": len(prs.slides)}
    payload.update(extra)
    print(json.dumps(payload, ensure_ascii=False))


# --------------------------------------------------------------------------- commands

def cmd_read(args) -> None:
    if args.markdown:
        try:
            from markitdown import MarkItDown

            print(MarkItDown().convert(str(args.file)).text_content)
            return
        except ImportError:
            print("markitdown not installed; falling back to python-pptx text extraction", file=sys.stderr)

    prs = _presentation(args.file)
    slides = []
    for index, slide in enumerate(prs.slides, start=1):
        title = _slide_title(slide)
        texts: list[str] = []
        for tf in _iter_text_frames(slide.shapes):
            for paragraph in tf.paragraphs:
                text = _paragraph_text(paragraph).strip()
                if text:
                    texts.append(text)
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        slides.append(
            {
                "index": index,
                "layout": slide.slide_layout.name,
                "title": title,
                "text": texts,
                "notes": notes,
            }
        )

    if args.json:
        print(json.dumps({"file": str(args.file), "slides": len(slides), "detail": slides}, ensure_ascii=False, indent=2))
        return

    for s in slides:
        print(f"--- Slide {s['index']} [{s['layout']}] {s['title']}".rstrip())
        for text in s["text"]:
            if text != s["title"]:
                print(f"  {text}")
        if s["notes"]:
            print(f"  [notes] {s['notes']}")
    if not slides:
        print("(deck has no slides)")


def cmd_layouts(args) -> None:
    prs = _presentation(args.file)
    for index, layout in enumerate(prs.slide_layouts):
        kinds = []
        for ph in layout.placeholders:
            kinds.append(str(ph.placeholder_format.type).split(".")[-1].split(" ")[0])
        print(f"  [{index}] {layout.name}  placeholders: {', '.join(kinds) or '-'}")


def _parse_slides(spec_json: str | None) -> list[dict]:
    if not spec_json:
        return []
    data = json.loads(spec_json)
    if not isinstance(data, list):
        raise SystemExit("Error: --slides-json must be a JSON array")
    slides = []
    for entry in data:
        if isinstance(entry, str):
            slides.append({"title": entry})
        elif isinstance(entry, dict):
            slides.append(entry)
        else:
            raise SystemExit(f"Error: unsupported slide entry {entry!r}")
    return slides


def cmd_create(args) -> None:
    prs = _presentation(args.template)
    if args.template and not args.keep_template_slides:
        for sld_id in list(prs.slides._sldIdLst):
            _drop_slide(prs, sld_id)

    title_layout = _pick_layout(prs, args.title_layout, default_names=("Title Slide",), default_index=0)
    title_slide = prs.slides.add_slide(title_layout)
    if title_slide.shapes.title is not None:
        title_slide.shapes.title.text = args.title
    if args.subtitle:
        for shape in title_slide.placeholders:
            if shape.placeholder_format.idx == 1 and shape.has_text_frame:
                shape.text_frame.text = args.subtitle
                break

    for spec in _parse_slides(args.slides_json):
        layout = _pick_layout(prs, spec.get("layout", args.layout))
        slide = prs.slides.add_slide(layout)
        _fill_slide(prs, slide, spec.get("title"), spec.get("bullets"), spec.get("notes"))

    out = Path(args.output)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    _summary(out, prs, created=True)


def cmd_add_slide(args) -> None:
    prs = _presentation(args.file)
    layout = _pick_layout(prs, args.layout)
    slide = prs.slides.add_slide(layout)
    bullets = json.loads(args.bullets_json) if args.bullets_json else []
    _fill_slide(prs, slide, args.title, bullets, args.notes)

    position = len(prs.slides)
    if args.index is not None:
        target = max(1, min(int(args.index), len(prs.slides)))
        sld_ids = prs.slides._sldIdLst
        element = sld_ids[-1]
        sld_ids.remove(element)
        sld_ids.insert(target - 1, element)
        position = target

    prs.save(str(args.file))
    _summary(args.file, prs, added_at=position, layout=layout.name)


def cmd_delete_slide(args) -> None:
    prs = _presentation(args.file)
    total = len(prs.slides)
    if not 1 <= args.index <= total:
        raise SystemExit(f"Error: slide index {args.index} out of range (1..{total})")
    sld_id = list(prs.slides._sldIdLst)[args.index - 1]
    _drop_slide(prs, sld_id)
    prs.save(str(args.file))
    _summary(args.file, prs, deleted=args.index)


def cmd_replace(args) -> None:
    prs = _presentation(args.file)
    count = 0
    slides_changed = 0
    for slide in prs.slides:
        before = count
        frames = list(_iter_text_frames(slide.shapes))
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame is not None:
            frames.append(slide.notes_slide.notes_text_frame)
        for tf in frames:
            for paragraph in tf.paragraphs:
                count += _replace_in_paragraph(paragraph, args.find, args.replace)
        if count != before:
            slides_changed += 1
    if count:
        prs.save(str(args.file))
    print(f"Replaced {count} occurrence(s) of {args.find!r} -> {args.replace!r} in {args.file}")
    print(json.dumps({"replaced": count, "slides_changed": slides_changed, "file": str(args.file)}))


def _libreoffice_cmd() -> str | None:
    for candidate in ("libreoffice", "soffice", "/Applications/LibreOffice.app/Contents/MacOS/soffice"):
        if "/" in candidate:
            if Path(candidate).exists():
                return candidate
            continue
        resolved = shutil.which(candidate)
        if resolved:
            return resolved
    return None


def cmd_to_pdf(args) -> None:
    src = Path(args.file).resolve()
    outdir = Path(args.outdir).resolve() if args.outdir else src.parent
    binary = _libreoffice_cmd()
    if not binary:
        raise SystemExit("Error: LibreOffice/soffice not found. Install LibreOffice to enable PPTX -> PDF conversion.")
    result = subprocess.run(
        [binary, "--headless", "--convert-to", "pdf", str(src), "--outdir", str(outdir)],
        capture_output=True,
        text=True,
    )
    out = outdir / src.with_suffix(".pdf").name
    if result.returncode != 0 or not out.exists():
        raise SystemExit(f"Error: LibreOffice conversion failed: {(result.stderr or result.stdout).strip()}")
    print(f"PDF: {out}")


# --------------------------------------------------------------------------- cli

def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    sub = parser.add_subparsers(dest="command", required=True)

    p = sub.add_parser("read", help="print slide text (and notes)")
    p.add_argument("file")
    p.add_argument("--json", action="store_true")
    p.add_argument("--markdown", action="store_true", help="use markitdown if installed")
    p.set_defaults(func=cmd_read)

    p = sub.add_parser("layouts", help="list slide layouts")
    p.add_argument("file")
    p.set_defaults(func=cmd_layouts)

    p = sub.add_parser("create", help="create a deck from a title + slide specs")
    p.add_argument("output")
    p.add_argument("--title", required=True)
    p.add_argument("--subtitle")
    p.add_argument("--slides-json")
    p.add_argument("--layout", help="default layout for content slides (index or name)")
    p.add_argument("--title-layout", help="layout for the title slide (index or name)")
    p.add_argument("--template", help="existing .pptx whose masters/layouts to use")
    p.add_argument("--keep-template-slides", action="store_true")
    p.set_defaults(func=cmd_create)

    p = sub.add_parser("add-slide", help="append (or insert) one slide")
    p.add_argument("file")
    p.add_argument("--title")
    p.add_argument("--bullets-json")
    p.add_argument("--notes")
    p.add_argument("--layout")
    p.add_argument("--index", type=int, help="1-based position to insert at")
    p.set_defaults(func=cmd_add_slide)

    p = sub.add_parser("delete-slide", help="remove one slide")
    p.add_argument("file")
    p.add_argument("--index", type=int, required=True)
    p.set_defaults(func=cmd_delete_slide)

    p = sub.add_parser("replace", help="find & replace text across slides, tables and notes")
    p.add_argument("file")
    p.add_argument("--find", required=True)
    p.add_argument("--replace", required=True)
    p.set_defaults(func=cmd_replace)

    p = sub.add_parser("to-pdf", help="convert via LibreOffice")
    p.add_argument("file")
    p.add_argument("--outdir")
    p.set_defaults(func=cmd_to_pdf)
    return parser


def main(argv: list[str] | None = None) -> None:
    args = build_parser().parse_args(argv)
    args.func(args)


if __name__ == "__main__":
    main()
